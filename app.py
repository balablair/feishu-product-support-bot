"""
飞书群消息监控机器人 - 长连接（WebSocket）模式
使用飞书官方 Python SDK lark-oapi，无需 Flask / ngrok / 公网地址
"""

import collections
import json
import logging
import os
import re
import threading
import time
from pathlib import Path
from threading import Lock

import requests
import lark_oapi as lark
import base64
from lark_oapi.api.im.v1 import (
    CreateMessageRequest,
    CreateMessageRequestBody,
    GetMessageResourceRequest,
    ReplyMessageRequest,
    ReplyMessageRequestBody,
)

from config import Config
from feedback import ai_detect_and_classify, save_feedback, should_reply
from docs import load_feishu_docs
from rag import RAGIndex

# ── 知识库加载 ────────────────────────────────────────────────────────────────
def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text.strip())
        if texts:
            slides.append(f"[第{i}页] " + " ".join(texts))
    return "\n".join(slides)


def _extract_pdf(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[第{i}页] {text.strip()}")
    return "\n".join(pages)


def load_knowledge() -> str:
    """
    读取 knowledge/ 目录下所有文件，合并为一段文本
    支持：.txt .md .docx .pptx .pdf
    启动时加载一次，之后常驻内存
    """
    kb_dir = Path(__file__).parent / "knowledge"
    if not kb_dir.exists():
        return ""

    extractors = {
        ".txt":  lambda f: f.read_text(encoding="utf-8"),
        ".md":   lambda f: f.read_text(encoding="utf-8"),
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".pdf":  _extract_pdf,
    }

    parts = []
    for f in sorted(kb_dir.glob("**/*")):
        if not f.is_file() or f.suffix not in extractors:
            continue
        try:
            text = extractors[f.suffix](f).strip()
            if text:
                parts.append(f"=== {f.name} ===\n{text}")
                # 日志在 logger 初始化后才能用，这里先跳过
        except Exception as e:
            # 解析失败不影响其他文件
            pass

    return "\n\n".join(parts)

# ── 日志配置 ──────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger(__name__)
# 屏蔽 pdfplumber 的字体解析噪音
logging.getLogger("pdfminer").setLevel(logging.ERROR)

# ── 启动时加载知识库 ──────────────────────────────────────────────────────────
KNOWLEDGE = load_knowledge()
if KNOWLEDGE:
    logger.info(f"本地知识库已加载，共 {len(KNOWLEDGE)} 字符")
else:
    logger.info("knowledge/ 目录为空，未加载本地知识库")

# ── 飞书 SDK 客户端（负责所有 API 调用）──────────────────────────────────────
feishu_client = (
    lark.Client.builder()
    .app_id(Config.FEISHU_APP_ID)
    .app_secret(Config.FEISHU_APP_SECRET)
    .build()
)

# ── 启动时加载飞书云文档知识库 ───────────────────────────────────────────────
_cloud_knowledge = load_feishu_docs(feishu_client)
if _cloud_knowledge:
    KNOWLEDGE = (KNOWLEDGE + "\n\n" + _cloud_knowledge).strip() if KNOWLEDGE else _cloud_knowledge
    logger.info(f"云文档知识库已合并，知识库总计 {len(KNOWLEDGE)} 字符")

# ── RAG 索引（向量化知识库，每次只检索相关片段）──────────────────────────────
_rag = RAGIndex(
    api_key=Config.EMBEDDING_API_KEY,
    embed_url=Config.EMBEDDING_API_URL,
    model=Config.EMBEDDING_MODEL,
    top_k=Config.RAG_TOP_K,
    chunk_size=Config.RAG_CHUNK_SIZE,
)
if KNOWLEDGE and Config.EMBEDDING_API_KEY:
    threading.Thread(target=_rag.build, args=(KNOWLEDGE,), daemon=True).start()
elif KNOWLEDGE:
    logger.warning("未配置 EMBEDDING_API_KEY，RAG 不可用，将使用全量知识库模式")

# ── 事件去重：避免飞书重复推送被重复处理 ────────────────────────────────────
_processed_events: dict[str, float] = {}
_event_lock = Lock()
EVENT_DEDUPE_TTL = 300  # 秒

# ── 对话历史（多轮上下文）────────────────────────────────────────────────────
# key: thread_id（thread 内对话）或 "chat:{chat_id}:user:{open_id}"（独立消息）
# value: deque of {"role": "user"/"assistant", "content": "..."}
_histories: dict[str, collections.deque] = {}
_history_timestamps: dict[str, float] = {}
_history_lock = Lock()
MAX_HISTORY_TURNS = Config.MAX_HISTORY_TURNS  # 每段对话保留的轮数（1轮 = user+assistant）
HISTORY_TTL = Config.HISTORY_TTL              # 无活动后多少秒重置对话（默认 30 分钟）


def _get_context_key(message, sender_id: str) -> str:
    """
    确定对话上下文 key：
    - 消息在 thread 中（root_id 不为空）→ 用 thread_id，thread 内所有人共享上下文
    - 普通群消息 → 用 chat+user，每个用户独立上下文
    """
    root_id = getattr(message, "root_id", "") or ""
    if root_id:
        return f"thread:{root_id}"
    return f"chat:{message.chat_id}:user:{sender_id}"


def _get_history(key: str) -> list[dict]:
    """获取对话历史，若已超过 TTL 则重置"""
    now = time.time()
    with _history_lock:
        if key in _history_timestamps and now - _history_timestamps[key] > HISTORY_TTL:
            _histories.pop(key, None)
            _history_timestamps.pop(key, None)
        return list(_histories.get(key, []))


def _add_to_history(key: str, user_text: str, assistant_reply: str) -> None:
    """将本轮对话追加到历史，超出 MAX_HISTORY_TURNS 时自动丢弃最旧的"""
    with _history_lock:
        if key not in _histories:
            # maxlen = 轮数 * 2（每轮 1 条 user + 1 条 assistant）
            _histories[key] = collections.deque(maxlen=MAX_HISTORY_TURNS * 2)
        _histories[key].append({"role": "user", "content": user_text})
        _histories[key].append({"role": "assistant", "content": assistant_reply})
        _history_timestamps[key] = time.time()


def _is_duplicate(event_id: str) -> bool:
    """基于 event_id 去重，TTL 内的相同事件视为重复"""
    now = time.time()
    with _event_lock:
        # 清理过期记录，防止内存无限增长
        expired = [k for k, t in _processed_events.items() if now - t > EVENT_DEDUPE_TTL]
        for k in expired:
            del _processed_events[k]

        if event_id in _processed_events:
            return True
        _processed_events[event_id] = now
        return False


# ── 飞书：回复消息（在 thread 中）────────────────────────────────────────────
def reply_to_message(message_id: str, text: str) -> bool:
    """
    回复指定消息，自动在 thread 中展开。
    首次回复会新建 thread，后续回复会追加到同一 thread。
    """
    try:
        req = (
            ReplyMessageRequest.builder()
            .message_id(message_id)
            .request_body(
                ReplyMessageRequestBody.builder()
                .content(json.dumps({"text": text}))
                .msg_type("text")
                .build()
            )
            .build()
        )
        resp = feishu_client.im.v1.message.reply(req)
        if not resp.success():
            logger.error(f"回复消息失败: code={resp.code}, msg={resp.msg}")
            return False
        logger.info(f"已回复消息 {message_id}")
        return True
    except Exception as e:
        logger.error(f"回复消息异常: {e}")
        return False


def send_text_message(chat_id: str, text: str) -> bool:
    """向指定群发送独立消息（不在 thread 中，用于主动通知场景）"""
    try:
        request = (
            CreateMessageRequest.builder()
            .receive_id_type("chat_id")
            .request_body(
                CreateMessageRequestBody.builder()
                .receive_id(chat_id)
                .msg_type("text")
                .content(json.dumps({"text": text}))
                .build()
            )
            .build()
        )
        resp = feishu_client.im.v1.message.create(request)
        if not resp.success():
            logger.error(f"发送消息失败: code={resp.code}, msg={resp.msg}")
            return False
        logger.info(f"消息已发送到群 {chat_id}")
        return True
    except Exception as e:
        logger.error(f"发送消息异常: {e}")
        return False


# ── AI：生成自然对话回复 ──────────────────────────────────────────────────────
def _load_soul() -> str:
    """加载 SOUL.md 作为 bot 的人格定义与行为规范"""
    soul_path = Path(__file__).parent / "SOUL.md"
    if soul_path.exists():
        text = soul_path.read_text(encoding="utf-8").strip()
        logger.info("SOUL.md 已加载")
        return text
    logger.warning("未找到 SOUL.md，将使用默认 prompt")
    return "你是一个飞书群助手，请简洁、友好地回复用户消息。"


# 启动时加载 SOUL.md
SOUL_PROMPT = _load_soul()


_NO_MARKDOWN = """
---
【格式要求（必须严格遵守）】
回复内容只能是纯文本。禁止使用任何 Markdown 语法，包括但不限于：
- 加粗（**文字** 或 __文字__）
- 斜体（*文字* 或 _文字_）
- 标题（# ## ###）
- 列表符号（- * 1.）
- 代码块（``` 或 `）
- 分隔线（---）
- 链接格式（[文字](url)）
如果需要列举内容，用"1、2、3、"或换行代替。"""


def _build_system_prompt(context: str = "") -> str:
    """
    组装最终 system prompt：SOUL.md + 知识库内容（RAG 检索片段或全量兜底）+ 格式要求
    context: RAG 检索出的相关片段；为空时退回全量知识库（RAG 未就绪时的兜底）
    """
    base = SOUL_PROMPT
    kb_content = context or KNOWLEDGE
    if kb_content:
        base += "\n\n---\n\n以下是与当前问题最相关的知识库内容，回答时优先依据此内容：\n\n" + kb_content
    return base + _NO_MARKDOWN


def generate_reply(text: str, history: list[dict] | None = None) -> str:
    """
    调用 AI API 生成回复。
    history 为本轮之前的对话记录（[{"role": "user/assistant", "content": "..."}]），
    传入后 AI 可感知对话上下文，实现多轮对话。
    失败时返回兜底回复，保证主流程不中断。
    """
    if not Config.OPENAI_API_KEY:
        logger.warning("未配置 OPENAI_API_KEY")
        return "收到，稍后回复你。"

    # 构建消息列表：system → 历史对话 → 当前用户消息
    context = _rag.retrieve(text) if _rag.is_ready else ""
    messages = [{"role": "system", "content": _build_system_prompt(context)}]
    if history:
        messages.extend(history)
    messages.append({"role": "user", "content": text})

    try:
        resp = requests.post(
            Config.OPENAI_API_URL,
            headers={
                "Authorization": f"Bearer {Config.OPENAI_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": Config.OPENAI_MODEL,
                "messages": messages,
                "temperature": 0.7,
            },
            timeout=15,
        )
        resp.raise_for_status()
        reply = resp.json()["choices"][0]["message"]["content"].strip()
        logger.info(f"AI 回复（历史 {len(history or [])} 条）: {reply[:100]}")
        return reply

    except Exception as e:
        logger.error(f"调用 AI API 失败: {e}")
        return "收到，稍后回复你。"


def download_image(message_id: str, image_key: str) -> str | None:
    """从飞书下载图片，返回 base64 字符串，失败返回 None"""
    try:
        req = (
            GetMessageResourceRequest.builder()
            .message_id(message_id)
            .file_key(image_key)
            .type("image")
            .build()
        )
        resp = feishu_client.im.v1.message_resource.get(req)
        if not resp.success():
            logger.error(f"图片下载失败: code={resp.code} msg={resp.msg}")
            return None
        return base64.b64encode(resp.file.read()).decode("utf-8")
    except Exception as e:
        logger.error(f"图片下载异常: {e}")
        return None


def generate_reply_with_image(image_b64: str, text: str = "") -> str:
    """调用 Qwen-VL 分析图片"""
    if not Config.VISION_API_KEY:
        return "收到，稍后回复你。"
    try:
        user_content = [
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_b64}"}},
            {"type": "text", "text": text if text else "请描述这张图片的内容，如果是报错截图请帮忙分析问题。"},
        ]
        resp = requests.post(
            Config.VISION_API_URL,
            headers={
                "Authorization": f"Bearer {Config.VISION_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": Config.VISION_MODEL,
                "messages": [
                    {"role": "system", "content": _build_system_prompt()},
                    {"role": "user", "content": user_content},
                ],
                "temperature": 0.7,
            },
            timeout=30,
        )
        resp.raise_for_status()
        reply = resp.json()["choices"][0]["message"]["content"].strip()
        logger.info(f"图片 AI 回复: {reply[:80]}")
        return reply
    except Exception as e:
        logger.error(f"图片 AI 调用失败: {e}")
        return "收到，稍后回复你。"


# ── 知识库热更新 ───────────────────────────────────────────────────────────────
def reload_knowledge() -> str:
    """重新加载本地文件和飞书云文档知识库，并异步重建 RAG 索引"""
    global KNOWLEDGE
    local = load_knowledge()
    cloud = load_feishu_docs(feishu_client)
    if local and cloud:
        KNOWLEDGE = local + "\n\n" + cloud
    else:
        KNOWLEDGE = local or cloud
    msg = f"知识库已重新加载，共 {len(KNOWLEDGE)} 字符。"
    logger.info(msg)
    if KNOWLEDGE and Config.EMBEDDING_API_KEY:
        def _rebuild():
            ok = _rag.rebuild(KNOWLEDGE)
            logger.info("RAG 索引重建完成" if ok else "RAG 索引重建失败，已回退到全量模式")
        threading.Thread(target=_rebuild, daemon=True).start()
        msg += " RAG 索引正在后台重建..."
    return msg


# ── 消息处理主逻辑 ────────────────────────────────────────────────────────────
def handle_message(data: lark.im.v1.P2ImMessageReceiveV1):
    """核心处理：提取文本/图片 → 带上下文调用 AI → thread 内回复"""
    try:
        message = data.event.message
        sender = data.event.sender

        # 忽略机器人自身发送的消息，避免循环回复
        if sender.sender_type == "app":
            logger.info("忽略机器人自身消息")
            return

        sender_id = sender.sender_id.open_id if sender.sender_id else "unknown"

        # 确定本条消息的对话上下文 key
        ctx_key = _get_context_key(message, sender_id)

        # ── 图片消息（纯图片 or 富文本 post 里包含图片）────────────────────────
        if message.message_type in ("image", "post"):
            content = json.loads(message.content)

            if message.message_type == "image":
                image_keys = [content.get("image_key", "")]
                text = ""
            else:
                image_keys = []
                text_parts = []
                for block in content.get("content", []):
                    for elem in block:
                        if elem.get("tag") == "img":
                            image_keys.append(elem.get("image_key", ""))
                        elif elem.get("tag") in ("text", "a"):
                            text_parts.append(elem.get("text", ""))
                text = re.sub(r"<at[^>]*>[^<]*</at>", "", " ".join(text_parts)).strip()

            image_keys = [k for k in image_keys if k]
            if not image_keys:
                if text:
                    logger.info(f"收到 post 文本消息 [{message.message_id}] 来自 {sender_id}: {text[:80]}")
                    history = _get_history(ctx_key)
                    reply = generate_reply(text, history)
                    reply_to_message(message.message_id, reply)
                    _add_to_history(ctx_key, text, reply)
                return

            logger.info(f"收到图片消息 [{message.message_id}] 来自 {sender_id}，图片数: {len(image_keys)}")
            image_b64 = download_image(message.message_id, image_keys[0])
            if image_b64:
                reply = generate_reply_with_image(image_b64, text)
                reply_to_message(message.message_id, reply)
                # 图片分析结果也写入历史，方便后续文字追问
                user_hist = f"[用户发送了图片]{(' ' + text) if text else ''}"
                _add_to_history(ctx_key, user_hist, reply)
                # 反馈检测：异步写入 Bitable
                feedback_text = f"{text}\n[图片内容：{reply}]" if text else f"[用户发送截图，图片内容：{reply}]"
                def _record_image(ft=feedback_text, ib=__import__('base64').b64decode(image_b64)):
                    category, summary = ai_detect_and_classify(ft)
                    if category:
                        save_feedback(feishu_client, sender_id, summary, message.chat_id, category, image_bytes=ib)
                threading.Thread(target=_record_image, daemon=True).start()
            else:
                reply_to_message(message.message_id, "图片下载失败，请重试。")
            return

        # ── 非文本且非图片，跳过 ──────────────────────────────────────────────
        if message.message_type != "text":
            logger.info(f"跳过不支持的消息类型: {message.message_type}")
            return

        # 解析消息文本，去掉 @mention 标记
        raw_text = json.loads(message.content).get("text", "").strip()
        text = re.sub(r"<at[^>]*>[^<]*</at>", "", raw_text).strip()
        was_mentioned = "<at" in raw_text

        logger.info(f"收到消息 [{message.message_id}] 来自 {sender_id}: {text[:80]}")

        # ── 管理员指令：/reload ────────────────────────────────────────────────
        if text.strip() == "/reload":
            if sender_id in Config.ADMIN_OPEN_IDS:
                status = reload_knowledge()
                reply_to_message(message.message_id, status)
            else:
                reply_to_message(message.message_id, "无权限执行此指令。")
            return

        # ── 过滤：@bot 直接回复；否则 AI 判断是否产品相关 ──────────────────
        if not was_mentioned and not should_reply(text):
            logger.info("消息与产品无关，跳过回复")
            return

        # ── 带上下文生成回复，在 thread 中回复 ───────────────────────────────
        history = _get_history(ctx_key)
        reply = generate_reply(text, history)
        reply_to_message(message.message_id, reply)
        _add_to_history(ctx_key, text, reply)

        # ── 反馈检测：AI 判断是否为反馈，是则异步写入多维表格 ────────────────
        def _record():
            category, summary = ai_detect_and_classify(text)
            if category:
                save_feedback(feishu_client, sender_id, summary, message.chat_id, category)
        threading.Thread(target=_record, daemon=True).start()

    except Exception as e:
        logger.error(f"处理消息时发生未预期错误: {e}", exc_info=True)


# ── 飞书 SDK 事件回调（注册到 EventDispatcher）────────────────────────────────
def on_message_receive(data: lark.im.v1.P2ImMessageReceiveV1) -> None:
    """
    SDK 在收到 im.message.receive_v1 事件时调用此函数
    SDK 要求 3 秒内返回，所以用子线程异步处理实际业务
    """
    event_id = data.header.event_id if data.header else ""

    # 去重：飞书可能对同一事件重复推送
    if event_id and _is_duplicate(event_id):
        logger.info(f"重复事件 {event_id}，跳过")
        return

    threading.Thread(target=handle_message, args=(data,), daemon=True).start()


# ── 程序入口 ──────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    # 配置检查
    try:
        Config.validate()
        logger.info("配置校验通过")
    except EnvironmentError as e:
        logger.warning(f"配置不完整: {e}，部分功能可能无法使用")

    # 注册事件处理器
    # builder 的两个参数（verify_token, encrypt_key）长连接模式下必须填空字符串
    event_handler = (
        lark.EventDispatcherHandler.builder("", "")
        .register_p2_im_message_receive_v1(on_message_receive)
        .build()
    )

    # 启动长连接 WebSocket 客户端，阻塞主线程
    logger.info("正在连接飞书长连接服务（WebSocket）...")
    ws_client = lark.ws.Client(
        Config.FEISHU_APP_ID,
        Config.FEISHU_APP_SECRET,
        event_handler=event_handler,
        log_level=lark.LogLevel.INFO,
    )
    ws_client.start()
